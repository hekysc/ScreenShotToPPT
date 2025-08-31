from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

def generate_ppt(folder, image_files, save_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for filename in image_files:
        path = os.path.join(folder, filename)
        img = Image.open(path)
        width, height = img.size

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        max_width = prs.slide_width * 0.9
        max_height = prs.slide_height * 0.9

        # print(width,max_width)
        # print(height,max_height)

        ratio = min(max_width / width, max_height / height)
        new_width = int(width * ratio)
        new_height = int(height * ratio)

        # print(ratio)

        # if width > max_width or height > max_height:
        #     ratio = min(max_width / width, max_height / height)
        #     new_width = int(width * ratio)
        #     new_height = int(height * ratio)
        # else:
        #     new_width = width
        #     new_height = height

        # print(prs.slide_width,new_width,prs.slide_height,new_height,filename)

        left = int((prs.slide_width - new_width) / 2)
        top = int((prs.slide_height - new_height) / 2)

        # print(left,top)

        slide.shapes.add_picture(path, left, top, width=new_width, height=new_height)

    prs.save(save_path)
